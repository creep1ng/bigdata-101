"""
Agrega slides de Continuous Pretraining y Destilación a la sesión 7
insertándolas después de LoRA/PEFT (slide 67, index 66) en la presentación.
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
import copy

prs = Presentation("PLN-sesiones-completas.pptx")
TITLE_CONTENT = prs.slide_layouts[1]

def make_slide(title, bullets):
    slide = prs.slides.add_slide(TITLE_CONTENT)
    slide.placeholders[0].text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(16)
        p.level = 1 if b.startswith("  ") else 0
        if b.startswith("  "):
            p.text = b.strip()
            p.font.size = Pt(14)
    return slide

# Create the 2 new slides (they get appended at the end)
make_slide("Continuous Pretraining", [
    "¿Qué es? Seguir entrenando un modelo base con datos de un dominio específico",
    "Diferencia con fine-tuning:",
    "  Fine-tuning: entrena para una TAREA (clasificar, responder)",
    "  Continuous pretraining: entrena para un DOMINIO (medicina, derecho, finanzas)",
    "  Usa la misma tarea del pre-training (predecir siguiente token / MLM)",
    "Ejemplo: tomar LLaMA y seguir entrenándolo con papers médicos",
    "  → El modelo aprende vocabulario, relaciones y conocimiento del dominio",
    "Resultado: un modelo base especializado, listo para fine-tuning posterior",
    "Casos reales: BioGPT (biomedicina), FinGPT (finanzas), CodeLLaMA (código)",
])

make_slide("Model Distillation (Destilación de modelos)", [
    "Problema: modelos grandes (GPT-4, LLaMA 70B) son costosos y lentos en producción",
    "Idea: transferir el conocimiento de un modelo GRANDE (teacher) a uno PEQUEÑO (student)",
    "El student aprende a imitar las probabilidades de salida del teacher",
    "  No solo la respuesta correcta, sino la distribución completa de probabilidades",
    "  Ejemplo: teacher dice P('gato')=0.7, P('felino')=0.2, P('perro')=0.05",
    "  El student aprende esas relaciones suaves entre palabras",
    "Resultado: modelo 3-10x más pequeño con 95-97% del rendimiento",
    "Casos reales: DistilBERT (60% tamaño, 97% rendimiento), TinyLLaMA",
    "Ideal para: producción, dispositivos móviles, bajo costo de inferencia",
])

# Now reorder: move the 2 new slides (last 2) to after slide index 66 (LoRA/PEFT)
sldIdLst = prs.slides._sldIdLst
sldId_list = list(sldIdLst)
# New slides are the last 2
new1 = sldId_list[-2]
new2 = sldId_list[-1]
# Remove them from end
sldIdLst.remove(new1)
sldIdLst.remove(new2)
# Insert after index 66 (slide 67 = LoRA/PEFT)
# Re-read the list after removal
sldId_list = list(sldIdLst)
target = sldId_list[66]  # LoRA/PEFT slide
target.addnext(new2)
target.addnext(new1)

prs.save("PLN-sesiones-completas.pptx")
print(f"✅ Guardado. Total slides: {len(prs.slides)}")

# Verify order around insertion
for i, slide in enumerate(prs.slides):
    if 65 <= i <= 71:
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None and shape.placeholder_format.idx == 0:
                title = shape.text_frame.text[:70]
                break
        print(f"  Slide {i+1}: {title}")
