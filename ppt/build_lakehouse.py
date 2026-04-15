"""
Genera la presentación de Data Lakes a Lakehouse usando template-upb.pptx.
Layouts: [0] Diapositiva de título, [1] Título y objetos, [2] Encabezado de sección
"""
from pptx import Presentation
from pptx.util import Pt

prs = Presentation("resources/template-upb.pptx")
TITLE_SLIDE = prs.slide_layouts[0]     # Portada
TITLE_CONTENT = prs.slide_layouts[1]   # Título y objetos
SECTION_HEADER = prs.slide_layouts[2]  # Encabezado de sección


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(TITLE_SLIDE)
    slide.placeholders[0].text = title
    slide.placeholders[1].text = subtitle


def add_section_slide(title, subtitle):
    slide = prs.slides.add_slide(SECTION_HEADER)
    slide.placeholders[0].text = title
    slide.placeholders[1].text = subtitle


def add_slide(title, bullets):
    slide = prs.slides.add_slide(TITLE_CONTENT)
    slide.placeholders[0].text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 1 if b.startswith("  ") else 0
        if b.startswith("  "):
            p.text = b.strip()
            p.font.size = Pt(14)
        else:
            p.text = b
            p.font.size = Pt(16)


# Eliminar la slide de ejemplo que trae el template
sldIdLst = prs.slides._sldIdLst
first = list(sldIdLst)[0]
sldIdLst.remove(first)
rId = first.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
prs.part.drop_rel(rId)

# ============================================================
# PORTADA
# ============================================================
add_title_slide(
    "De Data Lakes a Lakehouse",
    "Evolución de las arquitecturas de datos y el patrón Medallion\n"
    "Camilo Soto, MSc — Ingeniería en Ciencia de Datos — UPB"
)

# ============================================================
# AGENDA
# ============================================================
add_slide("Agenda", [
    "1. Data Warehouses: el punto de partida",
    "2. Data Lakes: la promesa del almacenamiento masivo",
    "3. Problemas del Data Lake: el Data Swamp",
    "4. Data Lakehouse: lo mejor de ambos mundos",
    "5. Arquitectura Medallion (Bronze, Silver, Gold)",
    "6. Catálogo y gobierno de datos (Unity Catalog)",
    "7. Tecnologías, herramientas y evolución histórica",
])

# ============================================================
# CAPÍTULO 1 — DATA WAREHOUSE
# ============================================================
add_section_slide("01 — Data Warehouse", "El modelo tradicional de análisis de datos")

add_slide("Data Warehouse — ¿Qué es?", [
    "Repositorio centralizado de datos estructurados",
    "Esquema definido antes de cargar → Schema-on-Write",
    "Optimizado para consultas analíticas (OLAP)",
    "Datos limpios, transformados y confiables",
    "SQL como lenguaje principal de consulta",
    "Ejemplos: Teradata, Oracle DW, Amazon Redshift, Snowflake (modo clásico)",
])

add_slide("Data Warehouse — Limitaciones", [
    "Solo datos estructurados (tablas relacionales)",
    "Costoso de escalar — almacenamiento propietario",
    "ETL rígido y lento ante cambios de esquema",
    "No soporta datos semi-estructurados ni no-estructurados",
    "  JSON anidado, logs, imágenes, audio → no caben",
    "Difícil integrar Machine Learning y Data Science",
    "  Los científicos de datos necesitan acceso a datos crudos",
    "Resultado: un sistema confiable pero limitado y caro",
])

# ============================================================
# CAPÍTULO 2 — DATA LAKE
# ============================================================
add_section_slide("02 — Data Lake", "Almacenamiento masivo y flexible")

add_slide("Data Lake — ¿Qué es?", [
    "Almacén de datos en formato crudo (raw), sin transformar",
    "Esquema al leer → Schema-on-Read",
    "Soporta cualquier tipo de dato:",
    "  CSV, JSON, Parquet, imágenes, audio, logs, video",
    "Almacenamiento barato y escalable:",
    "  Amazon S3, Azure Data Lake Storage (ADLS), HDFS",
    "Base para pipelines de Machine Learning y Big Data",
    "Idea: guardemos TODO, ya después vemos qué hacemos",
])

add_slide("Data Lake — Tecnologías clave", [
    "Almacenamiento: HDFS, Amazon S3, Azure ADLS, Google Cloud Storage",
    "Procesamiento: Apache Spark, Apache Flink, Presto, Trino",
    "Formatos de archivo:",
    "  Parquet — columnar, comprimido, estándar de facto",
    "  ORC — optimizado para Hive",
    "  Avro — orientado a filas, bueno para streaming",
    "  Delta, Iceberg, Hudi — formatos de tabla (veremos más adelante)",
    "Catálogo: Apache Hive Metastore, AWS Glue Catalog",
    "Consulta: Athena, Presto, Spark SQL, Hive",
])

# ============================================================
# CAPÍTULO 3 — DATA SWAMP
# ============================================================
add_section_slide("03 — Data Swamp", "Cuando el Data Lake se convierte en un pantano")

add_slide("Data Lake → Data Swamp — Problemas", [
    "Sin gobernanza: nadie sabe qué datos hay ni de dónde vienen",
    "Sin calidad: datos duplicados, incompletos, inconsistentes",
    "Sin transacciones ACID:",
    "  Si un job falla a mitad de escritura → datos corruptos",
    "  No hay rollback, no hay consistencia garantizada",
    "Sin versionamiento ni auditoría",
    "Rendimiento pobre en consultas ad-hoc",
    "Analogía: biblioteca sin catálogo donde cualquiera deja libros",
])

add_slide("Data Swamp — Consecuencias", [
    "Los analistas dejan de confiar en los datos del lake",
    "Se duplican esfuerzos: copias del lake a warehouses para poder usarlos",
    "  → Lo peor de ambos mundos: doble costo, doble complejidad",
    "Arquitecturas Lambda/Kappa para batch + streaming → muy complejas",
    "Costos ocultos de mantenimiento y debugging",
    "El lago se vuelve un pantano inutilizable",
    "Pregunta clave: ¿podemos tener lo mejor de ambos mundos?",
])

# ============================================================
# CAPÍTULO 4 — DATA LAKEHOUSE
# ============================================================
add_section_slide("04 — Data Lakehouse", "Lo mejor del Warehouse + lo mejor del Lake")

add_slide("Lakehouse — La convergencia", [
    "Antes: dos mundos separados",
    "  Warehouse → datos estructurados, SQL, BI, confiable pero caro",
    "  Lake → todos los datos, barato, flexible, pero caótico",
    "Ahora: un solo sistema unificado — el Lakehouse",
    "  Transacciones ACID sobre el lake",
    "  Schema enforcement + schema evolution",
    "  Soporte SQL + ML + Streaming en un solo lugar",
    "  Almacenamiento abierto y barato (S3, ADLS)",
    "  Gobernanza y auditoría integrada",
])

add_slide("Lakehouse — ¿Cómo funciona?", [
    "Capa de metadatos transaccional sobre archivos Parquet",
    "Transaction log: registro de cada operación (insert, update, delete)",
    "  Si un job falla → rollback automático",
    "Time travel: consultar datos como estaban en cualquier momento",
    "  SELECT * FROM tabla TIMESTAMP AS OF '2025-01-15'",
    "Schema enforcement: rechaza datos que no cumplen el esquema",
    "Schema evolution: agregar columnas sin romper nada",
    "El Lakehouse NO es un producto — es un patrón arquitectónico",
])

add_slide("Formatos de tabla abiertos", [
    "Delta Lake (Databricks, 2019)",
    "  ACID, time travel, schema evolution, optimización automática",
    "  Transaction log en JSON sobre Parquet",
    "Apache Iceberg (Netflix, 2018)",
    "  Partitioning oculto — el usuario no necesita saber la partición",
    "  Snapshots inmutables, evolución de esquema y partición",
    "Apache Hudi (Uber, 2016)",
    "  Upserts eficientes, Change Data Capture (CDC) incremental",
    "  Ideal para datos que cambian frecuentemente",
    "Los tres: open source, sobre Parquet, compatibles con Spark/Flink/Trino",
])

add_slide("Plataformas Lakehouse", [
    "Databricks Lakehouse Platform — líder, creador de Delta Lake",
    "AWS Lake Formation + Athena + Glue — serverless, integrado con S3",
    "Azure Synapse Analytics — integrado con ADLS y Power BI",
    "Google BigLake — unifica BigQuery con data lakes en GCS",
    "Snowflake (Iceberg Tables) — warehouse que adopta formato abierto",
    "Dremio — open source, motor SQL sobre Iceberg",
    "Apache Spark + Delta/Iceberg — DIY lakehouse",
])

# ============================================================
# CAPÍTULO 5 — ARQUITECTURA MEDALLION
# ============================================================
add_section_slide("05 — Arquitectura Medallion", "Bronze → Silver → Gold")

add_slide("El mundo ANTES de Medallion", [
    "Tenemos un Lakehouse con Delta Lake, ACID, time travel... genial.",
    "Pero ¿cómo organizamos los datos DENTRO del lake?",
    "Sin un patrón claro, cada equipo hacía lo que quería:",
    "  /datos/ventas_v2_final_FINAL.parquet",
    "  /raw/export_2024/ → ¿está limpio? ¿es confiable?",
    "  /analytics/reporte_juan/ → ¿de dónde sacó estos datos?",
    "Problemas concretos:",
    "  No hay separación entre datos crudos y datos listos para consumo",
    "  Si un pipeline falla, no hay forma clara de reprocesar",
    "  Nadie sabe qué datos están limpios y cuáles no",
    "Necesitamos un ESTÁNDAR de organización → Medallion",
])

add_slide("¿Qué es la arquitectura Medallion?", [
    "Patrón de diseño para organizar datos en un Lakehouse",
    "Popularizado por Databricks, hoy es estándar de la industria",
    "Idea central: los datos pasan por capas de calidad creciente",
    "  Como medallas olímpicas: Bronze → Silver → Gold",
    "  Cada capa tiene un propósito, un dueño y un nivel de confianza",
    "Bronze = datos crudos, tal cual llegan (materia prima)",
    "Silver = datos limpios, validados (producto en proceso)",
    "Gold = datos agregados, listos para consumo (producto terminado)",
    "Analogía: fábrica con almacén → control de calidad → vitrina",
])

add_slide("¿Qué resuelve Medallion?", [
    "1. Organización clara: cada dato tiene un lugar definido",
    "  No más /datos/ventas_v2_final_FINAL.parquet",
    "2. Reprocesamiento: si algo falla en Gold, reprocesas desde Bronze",
    "  Bronze es inmutable — siempre puedes volver al dato original",
    "3. Confianza: los consumidores saben que Gold está limpio y validado",
    "  No necesitan preguntarle al ingeniero '¿puedo confiar en esta tabla?'",
    "4. Separación de responsabilidades:",
    "  Ingenieros → Bronze y Silver | Analistas → Silver y Gold",
    "5. Linaje trazable: Bronze → Silver → Gold = origen claro de cada dato",
])

add_slide("Bronze — Capa de ingesta (Raw)", [
    "Ingesta directa sin transformar el CONTENIDO — los datos llegan tal cual",
    "Las fuentes pueden ser cualquier formato: CSV, JSON, XML, APIs, JDBC...",
    "  Se LEEN en su formato original (spark.read.csv, .json, .jdbc, etc.)",
    "  Se ESCRIBEN en formato Delta para tener ACID y time travel",
    "  El contenido no se modifica — solo cambia el formato de almacenamiento",
    "Append-only, inmutable — nunca se borran ni sobrescriben datos de Bronze",
    "Fuente de verdad histórica — si algo falla, reprocesas desde aquí",
    "Se agregan metadatos de ingesta: timestamp, archivo fuente, batch ID",
    "Tip: leer todo como STRING evita perder datos por errores de tipo",
])

add_slide("Silver — Capa de limpieza (Cleaned)", [
    "Datos limpiados, validados y conformados",
    "Deduplicación: eliminar registros duplicados",
    "Validación de tipos: strings a fechas, enteros, etc.",
    "Filtrado: remover registros inválidos o corruptos",
    "Joins entre fuentes: combinar datos de diferentes orígenes",
    "Schema enforcement: si un campo no cumple el tipo → tabla de errores",
    "Los datos de Silver son confiables para análisis y data science",
    "Ejemplo: users_cleaned, transactions_validated",
])

add_slide("Gold — Capa de negocio (Business)", [
    "Agregaciones de negocio: revenue mensual, usuarios activos, conversión",
    "KPIs y métricas calculadas según reglas del negocio",
    "Tablas dimensionales en star schema (hechos + dimensiones)",
    "Optimizado para dashboards de BI y consumo por ejecutivos",
    "Feature stores: variables calculadas para modelos de ML",
    "Acceso vía SQL, APIs, o herramientas de BI (Power BI, Tableau)",
    "Ejemplo: revenue_monthly, active_users_daily, churn_features",
])

add_slide("Medallion — Comparación por capa", [
    "Bronze | Calidad: cruda | Esquema: flexible | Usuarios: ingenieros",
    "  Actualización: append | Formato: Parquet/JSON | Ej: logs_raw",
    "Silver | Calidad: limpia | Esquema: enforced | Usuarios: analistas/DS",
    "  Actualización: merge/upsert | Formato: Delta/Iceberg | Ej: users_cleaned",
    "Gold | Calidad: agregada | Esquema: star schema | Usuarios: BI/ejecutivos",
    "  Actualización: scheduled | Formato: Delta/Iceberg | Ej: revenue_monthly",
])

add_slide("Beneficios de Medallion + Lakehouse", [
    "Para el equipo de datos:",
    "  Organización clara del pipeline de datos",
    "  Reprocesamiento fácil desde Bronze si algo falla",
    "  Linaje de datos trazable — sabes de dónde viene cada dato",
    "  Testing entre capas, separación de responsabilidades",
    "Para el negocio:",
    "  Datos confiables y auditables",
    "  Menor time-to-insight (sin copiar entre sistemas)",
    "  Un solo sistema para BI + ML",
    "  Costos reducidos vs warehouse separado",
])

add_slide("Medallion — Limitaciones y críticas", [
    "Duplicación de datos: Bronze + Silver + Gold = 3 copias del dato",
    "  En datasets masivos, el costo de storage se multiplica",
    "Latencia: flujo secuencial Bronze→Silver→Gold agrega delay",
    "  Complicado para casos de uso real-time / streaming",
    "Ingeniería pesada en Silver: deduplicar, limpiar, conformar",
    "  Requiere ingenieros especializados y mucho mantenimiento",
    "Falta de contexto de negocio en capas tempranas:",
    "  Bronze y Silver son técnicas — el negocio solo aparece en Gold",
    "  Los consumidores deben esperar a que Gold esté listo",
    "Medallion no es perfecto — es un punto de partida sólido",
])

add_slide("Más allá de Medallion — Enfoques complementarios", [
    "Data Mesh (Zhamak Dehghani, 2019):",
    "  Propiedad descentralizada — cada dominio es dueño de sus datos",
    "  No reemplaza Medallion, opera a otro nivel (organización, no pipeline)",
    "  Puedes tener Medallion DENTRO de cada dominio de un Data Mesh",
    "Data Vault (Dan Linstedt):",
    "  Modelado para warehouses con foco en historicidad y auditoría",
    "  Se puede usar dentro de Silver como técnica de modelado",
    "Data Fabric (Gartner):",
    "  Capa de integración con metadata activa y AI para conectar fuentes",
    "  Más un concepto de conectividad que de organización de capas",
    "Capa Platinum: extensión de Medallion para real-time y AI/ML",
])

add_slide("¿Cuándo usar qué?", [
    "Medallion: equipo centralizado, pipeline claro, batch/near-real-time",
    "  → La mayoría de empresas medianas. Punto de partida recomendado.",
    "Data Mesh + Medallion: organización grande con múltiples dominios",
    "  → Cada dominio tiene su propio Bronze/Silver/Gold",
    "Data Vault + Medallion: regulación fuerte, auditoría histórica",
    "  → Data Vault en Silver, Medallion como estructura general",
    "Streaming puro (Kappa): todo en tiempo real, sin batch",
    "  → IoT, fraude, trading. Medallion no aplica bien aquí.",
    "En la práctica: la mayoría combina patrones según la necesidad",
    "  No hay una arquitectura que 'supere' a todas las demás",
])

# ============================================================
# CAPÍTULO 6 — CATÁLOGO Y GOBIERNO DE DATOS
# ============================================================
add_section_slide("06 — Catálogo y Gobierno de Datos",
    "Unity Catalog: gobernanza centralizada para el Lakehouse")

add_slide("El mundo ANTES de Unity Catalog", [
    "Databricks usaba Hive Metastore — un catálogo por workspace",
    "Problema 1 — Silos: cada workspace tenía su propio metastore",
    "  Equipo A no podía ver las tablas del Equipo B",
    "  Para compartir datos → copiar archivos o dar acceso al storage",
    "Problema 2 — Seguridad débil: permisos se manejaban a nivel de storage",
    "  No había GRANT/REVOKE sobre tablas. Todo o nada.",
    "Problema 3 — Sin linaje ni auditoría: nadie sabía quién accedía a qué",
    "Problema 4 — Namespace de 2 niveles: schema.tabla (sin catálogo)",
    "  Difícil separar dev/prod o dominios de negocio",
])

add_slide("¿Qué es Unity Catalog?", [
    "Capa de gobernanza centralizada de Databricks (2022, open source 2024)",
    "UN solo metastore compartido por TODOS los workspaces de la región",
    "  → Se acabaron los silos. Todos ven el mismo catálogo de datos.",
    "Gobierna no solo tablas, sino TODOS los activos:",
    "  Tablas, vistas, volúmenes (archivos), modelos de ML, funciones",
    "Namespace de 3 niveles: catálogo.esquema.tabla",
    "  Antes: schema.tabla → Ahora: medallion.silver.travel_times",
    "  El nivel extra permite separar dominios, ambientes, proyectos",
    "Analogía: Hive Metastore era una biblioteca por piso.",
    "  Unity Catalog es UNA biblioteca central con catálogo unificado.",
])

add_slide("¿Qué resuelve Unity Catalog?", [
    "1. Descubrimiento: buscar tablas, ver esquemas, explorar datos",
    "  '¿Existe una tabla de clientes?' → búsqueda en el catálogo",
    "2. Control de acceso: GRANT/REVOKE con SQL estándar (ANSI SQL)",
    "  GRANT SELECT ON TABLE medallion.gold.metrics TO analistas",
    "  Herencia: permisos en catálogo → esquemas → tablas",
    "3. Linaje automático: de dónde viene cada tabla, quién la consume",
    "  gold.city_metrics ← silver.travel_times ← bronze.raw",
    "4. Auditoría: log completo de quién accedió a qué y cuándo",
    "  Cumplimiento regulatorio: GDPR, HIPAA, Habeas Data",
])

add_slide("Unity Catalog — Control de acceso en Medallion", [
    "Principio de mínimo privilegio por capa:",
    "  Ingenieros de datos → Bronze + Silver + Gold (lectura y escritura)",
    "  Científicos de datos → Silver + Gold (lectura)",
    "  Analistas de BI → Solo Gold (lectura)",
    "  Ejecutivos → Solo Gold vía dashboards (ni siquiera SQL directo)",
    "Ejemplo en SQL:",
    "  GRANT USE CATALOG ON CATALOG medallion TO ingenieros",
    "  GRANT SELECT ON SCHEMA medallion.gold TO analistas",
    "  REVOKE ALL ON SCHEMA medallion.bronze FROM analistas",
    "Los permisos viajan con los datos — no dependen del workspace",
])

add_slide("Unity Catalog — Linaje, tags y auditoría", [
    "Data Lineage (automático, sin código):",
    "  Visualización gráfica: tabla origen → transformaciones → tabla destino",
    "  ¿Un dashboard muestra datos raros? → seguir el linaje hasta Bronze",
    "Tags y clasificación de datos:",
    "  Marcar columnas sensibles: email → PII, salario → CONFIDENCIAL",
    "  Políticas automáticas: si tiene tag PII → solo acceso con aprobación",
    "Auditoría completa:",
    "  Cada SELECT, INSERT, UPDATE queda registrado",
    "  ¿Quién consultó la tabla de salarios el viernes a las 11pm?",
    "  Esencial para compliance y respuesta a incidentes",
])

add_slide("Unity Catalog — Estructura recomendada", [
    "Organización con Medallion + Unity Catalog:",
    "  Catálogo: medallion (o nombre del proyecto/dominio)",
    "  Esquemas: bronze, silver, gold",
    "  Tablas: medallion.bronze.travel_times_raw",
    "          medallion.silver.travel_times_clean",
    "          medallion.gold.city_metrics",
    "Volúmenes para archivos no tabulares:",
    "  CREATE VOLUME medallion.bronze.landing_files",
    "  Subir CSVs, PDFs, imágenes al volumen (reemplaza DBFS)",
    "Cada capa con permisos diferenciados por rol",
    "Esto es lo que van a ver en su workspace de Databricks",
])

# ============================================================
# EVOLUCIÓN HISTÓRICA Y CIERRE
# ============================================================
add_slide("Evolución histórica de las arquitecturas de datos", [
    "1990s — Data Warehouses (Teradata, Oracle, DB2)",
    "  Caros pero confiables. Solo datos estructurados.",
    "2006 — Hadoop y el nacimiento del Data Lake",
    "  Google publica MapReduce/GFS. Almacenamiento masivo en commodity HW.",
    "2010s — Cloud Lakes + Spark",
    "  S3/ADLS reemplazan HDFS. Spark reemplaza MapReduce.",
    "2020+ — Era del Lakehouse",
    "  Delta Lake, Iceberg, Hudi agregan la capa transaccional.",
    "  Medallion se convierte en el estándar de organización.",
    "Cada generación aprende de los errores de la anterior.",
])

add_slide("Resumen", [
    "1. Data Warehouse: confiable pero rígido y costoso",
    "2. Data Lake: flexible y barato pero caótico sin gobernanza",
    "3. Data Swamp: el fracaso del lake sin estructura",
    "4. Lakehouse: ACID + schema + SQL + ML sobre almacenamiento abierto",
    "5. Medallion: Bronze (raw) → Silver (clean) → Gold (business)",
    "6. Unity Catalog: gobernanza, linaje, control de acceso centralizado",
    "7. Tecnologías: Delta Lake, Iceberg, Hudi sobre Parquet",
    "El Lakehouse no es magia — es almacenamiento barato +",
    "  formatos abiertos + metadatos transaccionales + gobernanza",
])

add_section_slide("¿Preguntas?",
    "Camilo Soto, MSc\nIngeniería en Ciencia de Datos — UPB\nCátedra de Big Data"
)

# ============================================================
# GUARDAR
# ============================================================
prs.save("resources/lakehouse.pptx")
print(f"✅ Presentación guardada: resources/lakehouse.pptx")
print(f"   Total de diapositivas: {len(prs.slides)}")
