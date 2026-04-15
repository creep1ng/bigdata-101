"""
Agrega las diapositivas de las sesiones 5, 6 y 7 a la presentación PLN existente.
Usa layout [1] "Título y objetos" (título + contenido) y layout [0] para portadas de sección.
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
import copy

prs = Presentation("PLN-sesion 1 y 2.pptx")
TITLE_CONTENT = prs.slide_layouts[1]   # Título y objetos
SECTION_HEADER = prs.slide_layouts[2]  # Encabezado de sección

def add_section_slide(title, subtitle):
    slide = prs.slides.add_slide(SECTION_HEADER)
    slide.placeholders[0].text = title
    slide.placeholders[1].text = subtitle

def add_slide(title, bullets, sub_bullets=None):
    """Add a slide with title and bullet points. sub_bullets is a dict {index: [sub items]}"""
    slide = prs.slides.add_slide(TITLE_CONTENT)
    slide.placeholders[0].text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.level = 0
        if sub_bullets and i in sub_bullets:
            for sub in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.text = sub
                sp.font.size = Pt(14)
                sp.level = 1

# ============================================================
# SESIÓN 5 — Mecanismo de Atención y Arquitectura Transformer
# ============================================================

add_section_slide(
    "Sesión 5",
    "Del procesamiento secuencial a la atención: el Transformer"
)

add_slide("¿Dónde estamos? — Recapitulación", [
    "Sesiones 1-2: BoW, TF-IDF, n-grams → sin semántica, vectores dispersos",
    "Sesión 3: Word2Vec, LSTM, ELMo → secuencial, lento, dependencias largas",
    "Sesión 4: Evaluación de embeddings → necesitamos mejores representaciones",
    "Hoy: ¿Y si pudiéramos mirar toda la oración al mismo tiempo?",
])

add_slide("El cuello de botella de las LSTM", [
    "🔄 Secuencial → no paralelizable, entrenamiento lento",
    "📦 Compresión → toda la oración en un solo vector fijo",
    "📏 Distancia → información del inicio se degrada en oraciones largas",
    'Ejemplo: "El banco que está junto al río donde pescamos el verano pasado cerró sus puertas"',
    "La LSTM debe recorrer 14 palabras para conectar 'banco' con 'cerró'",
])

add_slide("Atención: la intuición", [
    "Analogía: cuando lees 'él', tu cerebro salta al sustantivo referido. No procesas linealmente.",
    "Atención = ponderar la importancia de cada palabra respecto a las demás",
    "Origen: Bahdanau et al. (2014) para traducción automática",
    "En vez de un vector fijo, el decoder 'mira' diferentes partes de la entrada",
    "contexto_t = Σ (α_t,i × h_i)",
    "Pero seguía usando LSTM → seguía siendo secuencial...",
])

add_slide('"Attention Is All You Need" (Vaswani et al., 2017)', [
    "El paper que cambió todo — Google, NeurIPS 2017",
    "Idea radical: eliminar completamente las LSTM. Usar SOLO atención.",
    "Antes: LSTM + atención como complemento → secuencial, meses de entrenamiento",
    "Después: Atención como arquitectura completa → paralelo, días de entrenamiento",
    "Resultado: nuevo estado del arte en traducción con mucho menos tiempo",
    "+130,000 citas. Origen de ChatGPT, Claude, Gemini, LLaMA...",
])

add_slide("Self-Attention: Query, Key, Value", [
    "Para cada palabra se calculan 3 vectores:",
    "Query (Q): ¿Qué estoy buscando?",
    "Key (K): ¿Qué información ofrezco?",
    "Value (V): ¿Cuál es mi contenido?",
    "Fórmula: Attention(Q, K, V) = softmax(Q × Kᵀ / √d_k) × V",
    "Q × Kᵀ → similitud | √d_k → escala | softmax → pesos | × V → combinación",
])

add_slide("Ejemplo: Self-Attention en acción", [
    '"El gato se sentó en la alfombra porque estaba cansado"',
    "Al procesar 'estaba', el modelo calcula pesos sobre todas las palabras:",
    "gato → 0.35 (sujeto) ← peso alto",
    "cansado → 0.30 (predicado) ← peso alto",
    "sentó → 0.08 | estaba → 0.10 | resto → ~0.02 c/u",
    "✅ Resuelve la correferencia SIN recorrer secuencialmente",
])

add_slide("Multi-Head Attention", [
    "En vez de una sola atención → 8 o 16 cabezas en paralelo",
    "Cada cabeza captura un tipo diferente de relación:",
    "Cabeza 1: relaciones sintácticas (sujeto ↔ verbo)",
    "Cabeza 2: relaciones semánticas (sinónimos)",
    "Cabeza 3: relaciones de proximidad",
    "MultiHead = Concat(head_1, ..., head_h) × W_O",
    "Como 8 analistas leyendo el mismo texto, cada uno buscando patrones diferentes",
])

add_slide("Codificación Posicional (Positional Encoding)", [
    "Problema: Self-attention trata la entrada como conjunto, no secuencia",
    '"El gato come pescado" = "Pescado come el gato" ← ¡mismo resultado!',
    "Solución: sumar un vector de posición a cada embedding",
    "PE(pos, 2i) = sin(pos / 10000^(2i/d))",
    "PE(pos, 2i+1) = cos(pos / 10000^(2i/d))",
    "Permiten aprender posiciones relativas y generalizar a secuencias más largas",
    "Entrada final = Embedding + Positional Encoding",
])

add_slide("Arquitectura completa del Transformer", [
    "ENCODER (×6 capas):",
    "  Multi-Head Self-Attention + Residual + LayerNorm",
    "  Feed-Forward Network + Residual + LayerNorm",
    "DECODER (×6 capas):",
    "  Masked Multi-Head Self-Attention (no ve el futuro)",
    "  Cross-Attention (Q del decoder, K/V del encoder)",
    "  Feed-Forward Network + Residual + LayerNorm",
    "Paper original: d_model=512, 8 heads, ~65M parámetros",
])

add_slide("Masked Attention y Cross-Attention", [
    "Masked Self-Attention (Decoder):",
    "  Al generar palabra t, NO puede ver palabras futuras (t+1, t+2...)",
    "  Máscara triangular: -∞ en posiciones futuras → softmax las convierte en 0",
    "Cross-Attention (Encoder ↔ Decoder):",
    "  Q viene del decoder ('¿qué necesito?')",
    "  K, V vienen del encoder ('aquí está la entrada completa')",
    "  Ejemplo: al generar 'gato' en español, atiende a 'cat' en inglés",
])

add_slide("¿Por qué ganó el Transformer?", [
    "Velocidad: procesamiento paralelo vs secuencial",
    "Dependencias: acceso directo a cualquier posición vs degradación con distancia",
    "Escalabilidad: más datos + cómputo = mejor (predecible)",
    "Entrenamiento: días vs semanas/meses",
    "Insight clave: cada token ve directamente a todos los demás en UNA operación",
    "→ Permitió entrenar modelos MUCHO más grandes en MUCHO más datos",
])

add_slide("Del Transformer a los modelos modernos", [
    "Transformer (2017) se dividió en 3 familias:",
    "Solo Encoder → BERT, RoBERTa, DistilBERT → comprensión",
    "Solo Decoder → GPT, GPT-2/3/4, LLaMA, Mistral → generación",
    "Encoder-Decoder → T5, BART, mBART, FLAN-T5 → transformación seq2seq",
    "Línea temporal:",
    "2017 Transformer → 2018 BERT+GPT → 2019 GPT-2,T5 → 2020 GPT-3",
    "→ 2022 ChatGPT → 2023 GPT-4, LLaMA → 2024+ Agentes, multimodal",
])

add_slide("Sesión 5 — Resumen", [
    "1. Las LSTM tienen un cuello de botella secuencial",
    "2. La atención permite ponderar relevancia entre todas las palabras",
    '3. "Attention Is All You Need" eliminó las LSTM completamente',
    "4. Self-Attention + Multi-Head + Positional Encoding = Transformer",
    "5. Tres familias emergieron: encoder-only, decoder-only, encoder-decoder",
    "Próxima sesión: BERT vs GPT vs T5 — ¿cuándo usar cada uno?",
    "Ref: arxiv.org/abs/1706.03762 | jalammar.github.io/illustrated-transformer/",
])

# ============================================================
# SESIÓN 6 — BERT, GPT y T5
# ============================================================

add_section_slide(
    "Sesión 6",
    "BERT, GPT y T5: las tres familias del Transformer"
)

add_slide("Recapitulación: las 3 familias", [
    "Transformer (2017) se puede usar por partes:",
    "Encoder-Only → comprensión de texto",
    "Decoder-Only → generación de texto",
    "Encoder-Decoder → transformación de secuencias",
    "Sesión anterior: entendimos la arquitectura completa",
    "Hoy: ¿qué pasa si usamos solo una parte? ¿Para qué sirve cada familia?",
])

add_slide("Encoder-Only: la idea", [
    "Usa solo el encoder del Transformer (sin decoder)",
    "Procesa toda la secuencia bidireccionalmente",
    "Cada token atiende a todos los demás (izquierda Y derecha)",
    "Produce representaciones ricas de cada token",
    "NO genera texto nuevo — COMPRENDE texto existente",
    "Tareas ideales: clasificación, NER, análisis de sentimiento, QA extractivo",
])

add_slide("BERT (Google, 2018)", [
    "Bidirectional Encoder Representations from Transformers",
    "110M params (base) / 340M (large) | 12/24 capas encoder",
    "Entrenado en Wikipedia + BookCorpus (~3.3B palabras)",
    "Pre-entrenamiento con 2 tareas:",
    'MLM: enmascara 15% de tokens → predice palabras ocultas bidireccionalmente',
    '  "El [MASK] se sentó en la [MASK]" → "gato", "alfombra"',
    "NSP: dado un par de oraciones, ¿la segunda sigue a la primera?",
    "Innovación clave: bidireccionalidad (GPT-1 solo veía izq→der)",
])

add_slide("BERT: ejemplos prácticos", [
    "Clasificación de sentimiento:",
    '  [CLS] Esta película es increíble [SEP] → Vector [CLS] → Positivo (0.95)',
    "Extracción de respuestas (QA):",
    '  [CLS] ¿Dónde nació Einstein? [SEP] Einstein nació en Ulm, Alemania [SEP]',
    '  → Predice posición inicio y fin → "Ulm, Alemania"',
    "El token [CLS] captura la representación global de la secuencia",
    "Mismo modelo pre-entrenado, diferente capa final → diferentes tareas",
])

add_slide("Variantes de BERT", [
    "RoBERTa (Facebook, 2019): más datos, sin NSP → supera BERT en todo",
    "ALBERT (Google, 2019): comparte params entre capas → 18x menos params",
    "DistilBERT (HuggingFace, 2019): destilación → 60% tamaño, 97% rendimiento, 2x rápido",
    "ELECTRA (Google, 2020): discriminador en vez de generador → más eficiente",
    "DeBERTa (Microsoft, 2020): atención desacoplada → supera RoBERTa",
    "Patrón: la comunidad tomó BERT y lo optimizó en todas las direcciones",
])

add_slide("Decoder-Only: la idea", [
    "Usa solo el decoder del Transformer (sin encoder)",
    "Procesa de izquierda a derecha (autoregresivo)",
    "Cada token solo atiende a los tokens anteriores (masked self-attention)",
    "Genera texto token por token",
    'Entrenado para predecir la siguiente palabra: "El gato se sentó en la" → "alfombra"',
    "Tareas ideales: generación de texto, chat, código, razonamiento",
])

add_slide("GPT: la evolución (OpenAI)", [
    "GPT-1 (2018): 117M params, BookCorpus → pre-training + fine-tuning",
    "GPT-2 (2019): 1.5B params, WebText → zero-shot emergente",
    "GPT-3 (2020): 175B params, ~570GB texto → few-shot sin fine-tuning",
    "GPT-4 (2023): ~1.7T params (est.), multimodal → razonamiento avanzado",
    "Insight de escala: con suficientes params y datos, el modelo aprende",
    "tareas que nunca le enseñaron explícitamente (capacidades emergentes)",
])

add_slide("GPT: In-Context Learning", [
    "La revolución de GPT-3: aprender de ejemplos en el prompt",
    'Zero-shot: "Traduce al español: The cat is on the mat" → responde directo',
    'One-shot: "Hello → Hola. Goodbye →" → infiere "Adiós"',
    'Few-shot: "Me encanta→Positivo, Es horrible→Negativo, No está mal→" → "Positivo"',
    "Sin cambiar un solo parámetro del modelo. Solo con el contexto del prompt.",
    "Antes: miles de ejemplos + fine-tuning. Ahora: 2-3 ejemplos en el prompt.",
])

add_slide("Otros modelos Decoder-Only", [
    "LLaMA (Meta, 2023): 7B-70B params, open-source, eficiente",
    "LLaMA 2/3 (Meta, 2023-24): chat optimizado, licencia abierta",
    "Mistral 7B (Mistral AI, 2023): sliding window attention, muy eficiente",
    "Mixtral 8×7B (Mistral AI, 2024): Mixture of Experts, activa 2 de 8",
    "Gemma (Google, 2024): 2B-27B, ligero, open-source",
    "Qwen (Alibaba, 2024): 0.5B-72B, multilingüe",
    "Tendencia: modelos más pequeños, más eficientes, open-source",
])

add_slide("Encoder-Decoder: T5 y BART", [
    "Usa ambas partes del Transformer original",
    "Encoder procesa entrada (bidireccional) → Decoder genera salida (autoregresivo)",
    "T5 (Google, 2019): reformula TODA tarea como texto→texto",
    '  "classify: This movie is great" → "positive"',
    '  "translate English to German: Hello" → "Hallo"',
    '  "summarize: [artículo]" → "[resumen]"',
    "BART (Facebook, 2019): encoder bidireccional + decoder autoregresivo",
    "  Pre-entrenado corrompiendo texto y reconstruyéndolo",
])

add_slide("Comparación directa de las 3 familias", [
    "Encoder-Only (BERT): bidireccional | comprensión | predice masked tokens",
    "  → '¿Es positivo este review?' | BERT, RoBERTa, DeBERTa",
    "Decoder-Only (GPT): izq→der | generación | predice siguiente token",
    "  → 'Escríbeme un poema' | GPT-4, LLaMA, Mistral",
    "Encoder-Decoder (T5): bidireccional + izq→der | transformación seq2seq",
    "  → 'Resume este artículo' | T5, BART, FLAN-T5",
    "Tendencia 2024: decoder-only domina, pero BERT vive en producción (eficiencia)",
])

add_slide("¿Cuándo usar cada uno?", [
    "🔍 Entender/clasificar texto → Encoder-only (BERT, DeBERTa)",
    "  Sentimiento, NER, detección de spam, QA extractivo",
    "✍️ Generar texto → Decoder-only (GPT, LLaMA)",
    "  Chatbots, código, razonamiento, escritura creativa",
    "🔄 Transformar texto → Encoder-decoder (T5, BART)",
    "  Traducción, resumen, parafraseo, corrección gramatical",
    "Nota: los decoder-only modernos son tan grandes que también comprenden bien",
])

add_slide("Sesión 6 — Resumen", [
    "1. Encoder-only (BERT): bidireccional, ideal para comprensión",
    "2. Decoder-only (GPT): autoregresivo, ideal para generación",
    "3. Encoder-decoder (T5): ambos, ideal para transformación seq2seq",
    "4. La tendencia actual favorece decoder-only por su versatilidad",
    "5. Pero cada familia tiene su nicho en producción",
    "Próxima sesión: Fine-tuning, prompt engineering, vectorDB y RAG",
    "Trabajo 2: Modelo basado en Transformers para clasificación",
])

# ============================================================
# SESIÓN 7 — Fine-Tuning, Prompt Engineering, VectorDB, RAG
# ============================================================

add_section_slide(
    "Sesión 7",
    "Adaptación de modelos: Fine-Tuning, Prompt Engineering y RAG"
)

add_slide("El problema de la adaptación", [
    "Tenemos modelos pre-entrenados poderosos (BERT, GPT, LLaMA)",
    "Pero necesitamos que resuelvan NUESTRO problema específico",
    "Estrategia 1 — Fine-tuning: reentrenar con datos propios (modifica el modelo)",
    "Estrategia 2 — Prompt engineering: diseñar instrucciones inteligentes (no modifica)",
    "Problema adicional: los modelos tienen conocimiento limitado",
    "  Fecha de corte + sin datos privados",
    "Estrategia 3 — RAG: recuperar información externa para fundamentar respuestas",
])

add_slide("Fine-Tuning: concepto", [
    "Paradigma de transferencia de aprendizaje:",
    "Fase 1 — Pre-training (costoso, una vez):",
    "  Modelo aprende lenguaje general con millones/billones de textos",
    "Fase 2 — Fine-tuning (barato, muchas veces):",
    "  Modelo se adapta a tarea específica con datos etiquetados propios",
    "Analogía: Pre-training = universidad general. Fine-tuning = especialización.",
])

add_slide("Fine-Tuning clásico (Full)", [
    "Se reentrenan TODOS los parámetros del modelo",
    "BERT pre-entrenado (110M params) + capa de clasificación nueva",
    "+ Dataset etiquetado (ej: 10K reviews) + 3-5 epochs, lr=2e-5",
    "= Modelo especializado en sentimiento",
    "Ventajas: máximo rendimiento, adaptación completa",
    "Desventajas: necesita GPU, datos etiquetados, un modelo por tarea",
    "Con HuggingFace: from_pretrained() → Trainer → train() → listo",
])

add_slide("Fine-Tuning eficiente: LoRA y PEFT", [
    "Problema: fine-tuning completo de LLaMA 70B requiere ~140GB VRAM",
    "Solución: Parameter-Efficient Fine-Tuning (PEFT)",
    "LoRA: agrega matrices de bajo rango a capas de atención (~0.1% params)",
    "  W_original (congelado) + ΔW = A × B (entrenado), con rango r << d",
    "QLoRA: LoRA + cuantización a 4 bits → aún menos memoria",
    "→ Fine-tuning de LLaMA 7B en una GPU de 16GB con QLoRA",
    "Adapters (~1-5% params) | Prefix Tuning (~0.1% params)",
])

add_slide("Prompt Engineering: concepto", [
    "En vez de modificar el modelo, modificamos la ENTRADA",
    "El prompt es la instrucción/contexto que guía la respuesta",
    'Prompt malo: "sentimiento: me gusta esta película"',
    'Prompt bueno: "Eres un analista de sentimiento experto.',
    '  Clasifica como POSITIVO, NEGATIVO o NEUTRO.',
    '  Responde solo con la clasificación.',
    '  Texto: Me gusta esta película. Clasificación:"',
    "No cambia ningún parámetro. Solo cambia cómo le hablamos.",
])

add_slide("Técnicas de Prompt Engineering", [
    "Zero-shot: solo instrucción, sin ejemplos",
    "Few-shot: instrucción + ejemplos demostrativos",
    "Chain-of-Thought: pedir razonamiento paso a paso",
    '  Sin CoT: "¿15 × 27?" → "395" ❌',
    '  Con CoT: "Piensa paso a paso" → 15×20=300, 15×7=105, 300+105=405 ✅',
    "Role prompting: asignar un rol ('Eres un médico experto en...')",
    "Self-consistency: generar múltiples respuestas y votar la mayoría",
])

add_slide("Fine-Tuning vs Prompt Engineering", [
    "Fine-Tuning: miles de ejemplos, costo alto, rendimiento máximo, un modelo/tarea",
    "Prompt Eng: 0 a pocos ejemplos, costo bajo, flexible, iteración rápida",
    "Regla práctica:",
    "  Prototipo rápido → prompt engineering",
    "  Producción con alta precisión → fine-tuning",
    "  Modelo grande + datos privados → RAG",
    "Muchos proyectos: empiezan con prompts, migran a fine-tuning si necesitan más",
])

add_slide("El problema del conocimiento limitado", [
    "Los LLMs tienen 2 limitaciones críticas:",
    "1. Fecha de corte: no saben qué pasó después de su entrenamiento",
    "2. Sin datos privados: no conocen documentos internos de tu empresa",
    '"¿Cuál es la política de vacaciones de nuestra empresa?"',
    '→ "No tengo acceso a las políticas internas de su empresa."',
    "¿Cómo darle acceso a información externa sin reentrenarlo?",
    "→ Retrieval-Augmented Generation (RAG)",
])

add_slide("Bases de datos vectoriales", [
    "Almacenan embeddings y permiten buscar por similitud semántica",
    'Documento: "La política de vacaciones otorga 15 días al año"',
    "  → Modelo de embeddings → Vector [0.023, -0.156, 0.891, ...] (1536 dims)",
    "  → Se almacena en la base vectorial",
    'Búsqueda: "¿cuántos días de descanso tengo?" → encuentra el doc de vacaciones',
    "  Aunque no comparten palabras exactas, los vectores son cercanos",
    "Populares: ChromaDB, Pinecone, FAISS, Qdrant, pgvector",
])

add_slide("RAG: Retrieval-Augmented Generation", [
    "RAG = Recuperar información relevante + Generar respuesta fundamentada",
    "Flujo:",
    "1. Pregunta del usuario → embedding de la pregunta",
    "2. Búsqueda por similitud en base vectorial → top-k documentos",
    "3. Prompt = instrucción + documentos recuperados + pregunta",
    "4. LLM genera respuesta fundamentada en los documentos",
    "Como darle al modelo un 'cheat sheet' personalizado para cada pregunta",
])

add_slide("RAG: paso a paso", [
    "Fase de indexación (una vez):",
    "  Recopilar docs → dividir en chunks (~500-1000 tokens)",
    "  → generar embeddings → almacenar en base vectorial",
    "Fase de consulta (cada pregunta):",
    "  Embedding de pregunta → buscar k chunks similares",
    "  → construir prompt con contexto + pregunta → enviar al LLM",
    'Prompt típico: "Responde SOLO con info del contexto. Si no sabes, dilo."',
    "Chunking con overlap: cada fragmento comparte texto con el anterior/siguiente",
])

add_slide("RAG: ventajas y desafíos", [
    "Ventajas:",
    "  ✅ Acceso a información actualizada y privada",
    "  ✅ Respuestas fundamentadas, reducción de alucinaciones",
    "  ✅ Sin reentrenar el modelo. Fuentes citables y verificables",
    "Desafíos:",
    "  ⚠️ Calidad del chunking (muy grande pierde detalle, muy pequeño pierde contexto)",
    "  ⚠️ Calidad del modelo de embeddings y relevancia de docs recuperados",
    "  ⚠️ Ventana de contexto limitada del LLM. Latencia adicional",
])

add_slide("Comparación: Fine-Tuning vs Prompt Eng. vs RAG", [
    "Fine-Tuning: tarea específica, datos etiquetados, conocimiento en pesos, costo alto",
    "Prompt Eng: prototipo rápido, tareas generales, solo ventana de contexto, costo bajo",
    "RAG: datos privados/actualizados, docs externos, costo medio, menos alucinaciones",
    "Se pueden combinar: Fine-tuning + RAG + Prompt Engineering",
    "Ejemplo: LLaMA + LoRA (dominio) + ChromaDB (docs) + prompt (formato)",
])

add_slide("Ejemplo integrador: Chatbot de soporte técnico", [
    "1. Base de conocimiento: manuales → chunks → embeddings → ChromaDB",
    "2. Modelo: LLaMA 3 8B + LoRA fine-tuned con historial de tickets",
    "3. Cada consulta (RAG): pregunta → búsqueda ChromaDB → top-5 chunks",
    "   → prompt con contexto + pregunta → LLaMA → respuesta",
    '4. Prompt: "Eres un agente de soporte amable y preciso.',
    '   Responde SOLO con info del contexto.',
    '   Si no sabes, di: Voy a escalar tu caso."',
    "Fine-tuning + RAG + Prompt Engineering trabajando juntos",
])

add_slide("Sesión 7 — Resumen", [
    "1. Fine-tuning: adaptar el modelo reentrenando (full o LoRA)",
    "2. Prompt engineering: guiar el modelo con instrucciones inteligentes",
    "3. Bases vectoriales: almacenar y buscar por similitud semántica",
    "4. RAG: recuperar información relevante para fundamentar respuestas",
    "Próxima sesión: implementación práctica de RAG, ventana de contexto,",
    "  model distillation y continuous pretraining",
    "Trabajo 3: Implementación de un modelo RAG para ingeniería de prompts",
])

# Guardar
prs.save("PLN-sesiones-completas.pptx")
print(f"✅ Presentación guardada: PLN-sesiones-completas.pptx")
print(f"   Total de diapositivas: {len(prs.slides)}")
